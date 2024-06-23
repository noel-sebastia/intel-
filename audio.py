import os
import re
from flask import Flask, render_template, request, send_file
from bs4 import BeautifulSoup
import requests
from fpdf import FPDF
from pptx import Presentation
from docx import Document
from pptx.util import Inches
from transformers import pipeline, AutoModelForQuestionAnswering, AutoTokenizer
from celery import Celery
from pymongo import MongoClient
from datetime import datetime

# Initialize Flask app
app = Flask(__name__)

# Initialize Celery
app.config['CELERY_BROKER_URL'] = 'redis://localhost:6379/0'
app.config['CELERY_RESULT_BACKEND'] = 'redis://localhost:6379/0'
celery = Celery(app.name, broker=app.config['CELERY_BROKER_URL'])
celery.conf.update(app.config)

# Initialize MongoDB
client = MongoClient("mongodb://localhost:27017/")
db = client["scraped_data"]

# Initialize NLP pipeline for question answering (disabling GPU usage)
qa_model_name = "distilbert-base-cased-distilled-squad"
tokenizer = AutoTokenizer.from_pretrained(qa_model_name)
model = AutoModelForQuestionAnswering.from_pretrained(qa_model_name)
qa_pipeline = pipeline('question-answering', model=model, tokenizer=tokenizer, device=-1)

# Route for home page
@app.route('/')
def home():
    return render_template('scrape.html')

# Route for handling search request
@app.route('/search', methods=['POST'])
def search():
    query = request.form.get('query')
    query_type = request.form.get('query_type')
    
    if not query or not query_type:
        return render_template('scrpe.html', error="Please provide both query and query type.")
    
    keywords = extract_keywords(query)
    result = scrape_information(' '.join(keywords), query_type)
    
    return render_template('scrpe.html', query=query, query_type=query_type, result=result)

# Function to extract keywords from a query
def extract_keywords(query):
    return re.findall(r'\b\w+\b', query)

# Function to scrape information from different sources
def scrape_information(query, query_type):
    sources = [
        f"https://www.byjus.com/{query}",
        f"https://www.wikihow.com/{query}",
        f"https://simple.wikipedia.org/wiki/{query}",
        f"https://scholar.google.com/scholar?q={query}",
        f"https://en.wikipedia.org/wiki/{query}",
        f"https://www.britannica.com/{query}",
        f"https://sw.wikipedia.org/wiki/{query}"  # Adding Swahili Wikipedia
    ]
    content_list, images, related_keywords = [], [], []
    
    for url in sources:
        response = requests.get(url)
        if response.status_code == 200:
            soup = BeautifulSoup(response.content, 'html.parser')
            content_div = soup.find('div', class_='content') or soup.find('div', {'id': 'mw-content-text'})
            
            if content_div:
                paragraphs = content_div.find_all('p', limit=25)
                content_list.extend([p.text.strip() for p in paragraphs if p.text.strip()])
                
                for img in content_div.find_all('img', limit=5):
                    img_url = img.get('src')
                    if img_url:
                        images.append(f"https:{img_url}")
                
                if content_list:
                    related_keywords = extract_keywords(soup.title.string)
                    break
    
    if not content_list:
        return {"content": "Content not found.", "images": images, "related_keywords": related_keywords[:3]}
    
    summarized_content = summarize_content(content_list, query_type)
    save_to_db(query, summarized_content, images, related_keywords)
    
    return {"content": summarized_content, "images": images, "related_keywords": related_keywords[:3]}

# Function to summarize content based on query type
def summarize_content(content_list, query_type):
    if query_type == "definition":
        return content_list[0] if content_list else "Definition not found."
    elif query_type == "essay":
        return ' '.join(content_list[:300]) if content_list else "Essay content not found."
    elif query_type == "analysis":
        return ' '.join(content_list[:25]) if content_list else "Analysis not found."
    elif query_type == "description":
        return ' '.join(content_list[:10]) if content_list else "Description not found."
    else:
        return "Invalid query type."

# Function to save data to MongoDB
def save_to_db(query, content, images, related_keywords):
    db.data.insert_one({
        "query": query,
        "content": content,
        "images": images,
        "related_keywords": related_keywords,
        "timestamp": datetime.now()
    })

# Asynchronous task to update data in the background
@celery.task
def update_data_in_background(query, query_type):
    result = scrape_information(query, query_type)
    save_to_db(query, result['content'], result['images'], result['related_keywords'])

# Route for downloading files
@app.route('/download/<file_type>', methods=['POST'])
def download(file_type):
    query = request.form.get('query')
    query_type = request.form.get('query_type')
    content = request.form.get('content')
    
    if not query or not query_type or not content:
        return render_template('scrpe.html', error="Missing data to generate file.")
    
    result = {"content": content, "images": []}
    
    if file_type == 'pdf':
        create_pdf(query, query_type, result)
        return send_file("output.pdf", as_attachment=True)
    elif file_type == 'ppt':
        create_ppt(query, query_type, result)
        return send_file("Nobestudy.pptx", as_attachment=True)
    elif file_type == 'docx':
        create_docx(query, query_type, result)
        return send_file("Nobestudy.docx", as_attachment=True)
    else:
        return render_template('scrpe.html', error="Invalid file type.")

# Function to create PDF
def create_pdf(query, query_type, result):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    title_text = f"{query_type.capitalize()} of {query}"
    content_text = result['content']
    
    try:
        pdf.cell(200, 10, txt=title_text, ln=True, align='C')
        pdf.multi_cell(0, 10, txt=content_text)
        
        if result['images']:
            for img in result['images']:
                pdf.image(img, w=100)
        
        pdf.output("output.pdf")
    except Exception as e:
        print(f"An error occurred while creating the PDF: {e}")

# Function to create PowerPoint presentation
def create_ppt(query, query_type, result):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"{query_type.capitalize()} of {query}"
    
    left_inch = Inches(1)
    top_inch = Inches(1.2)
    width_inch = Inches(8)
    height_inch = Inches(1.5)
    
    textbox = slide.shapes.add_textbox(left_inch, top_inch, width_inch, height_inch)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = result['content']
    
    if result['images']:
        for img in result['images']:
            slide.shapes.add_picture(img, Inches(1), Inches(2), height=Inches(1.5))
    
    prs.save('Nobestudy.pptx')

# Function to create DOCX document
def create_docx(query, query_type, result):
    doc = Document()
    doc.add_heading(f"{query_type.capitalize()} of {query}", level=1)
    doc.add_paragraph(result['content'])
    
    if result['images']:
        for img in result['images']:
            doc.add_picture(img, width=Inches(2))
    
    doc.save('Nobestudy.docx')

# Run Flask app
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 4000))
    app.run(host="0.0.0.0", port=port)
