import os
import re
import random
from datetime import datetime
from flask import Flask, request, send_file, jsonify, render_template
from bs4 import BeautifulSoup
import requests
from fpdf import FPDF
from pptx import Presentation
from docx import Document
from pptx.util import Inches
from celery import Celery
from pymongo import MongoClient
from langdetect import detect
from googlesearch import search
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer
from googlesearch import search as google_search


# Uncomment to download NLTK resources if needed
# import nltk
# nltk.download('punkt')
# nltk.download('stopwords')

# Initialize Flask app
app = Flask(__name__)

# Initialize Celery
app.config['CELERY_BROKER_URL'] = 'redis://localhost:6379/0'
app.config['CELERY_RESULT_BACKEND'] = 'redis://localhost:6379/0'
celery = Celery(app.name, broker=app.config['CELERY_BROKER_URL'])
celery.conf.update(app.config)

# Initialize MongoDB
client = MongoClient("mongodb://localhost:27017/")
db = client["scraped_data"]

# List of domains to avoid in content extraction
avoid_domains = ["wikipedia.org", "byjus.com", "wikihow.com", "britannica.com"]

# Route for home page
@app.route('/')
def home():
    return render_template('scrpe.html')

# Route for handling search request
@app.route('/search', methods=['POST'])
def search():
    data = request.get_json()
    query = data.get('query')
    query_type = data.get('query_type')
    
    if not query or not query_type:
        return jsonify({"error": "Please provide both query and query type."}), 400
    
    keywords = extract_keywords(query)
    result = scrape_information(' '.join(keywords), query_type, query)
    
    return jsonify(result)

# Function to extract keywords from a query
def extract_keywords(query):
    return re.findall(r'\b\w+\b', query)

# Function to scrape information from different sources or Google if needed
def scrape_information(query, query_type, original_query):
    sources = [
        f"https://simple.wikipedia.org/wiki/{query}",
        f"https://scholar.google.com/scholar?q={query}",
        f"https://en.wikipedia.org/wiki/{query}",
        f"https://sw.wikipedia.org/wiki/{query}"
    ]
    content_list, images, related_keywords = [], [], []
    
    # Scraping from predefined sources
    for url in sources:
        if any(domain in url for domain in avoid_domains):
            continue
        
        response = requests.get(url)
        if response.status_code == 200:
            soup = BeautifulSoup(response.content, 'html.parser')
            content_div = soup.find('div', class_='content') or soup.find('div', {'id': 'mw-content-text'})
            
            if content_div:
                paragraphs = content_div.find_all('p', limit=25)
                content_list.extend([clean_text(p.text.strip()) for p in paragraphs if p.text.strip()])
                
                for img in content_div.find_all('img', limit=5):
                    img_url = img.get('src')
                    if img_url:
                        if img_url.startswith('//'):
                            img_url = 'https:' + img_url
                        images.append(img_url)
                
                if content_list:
                    related_keywords = extract_keywords(soup.title.string)
                    break
    
def scrape_information(query, query_type, original_query):
    sources = [
        f"https://www.byjus.com/{query}",
        f"https://www.wikihow.com/{query}",
        f"https://simple.wikipedia.org/wiki/{query}",
        f"https://scholar.google.com/scholar?q={query}",
        f"https://en.wikipedia.org/wiki/{query}",
        f"https://sw.wikipedia.org/wiki/{query}",
        f"https://www.britannica.com/{query}"
    ]
    content_list, images, related_keywords = [], [], []
    
    # Scrape content from predefined sources
    for url in sources:
        response = requests.get(url)
        if response.status_code == 200:
            soup = BeautifulSoup(response.content, 'html.parser')
            content_div = soup.find('div', class_='content') or soup.find('div', {'id': 'mw-content-text'})
            
            if content_div:
                paragraphs = content_div.find_all('p', limit=25)
                content_list.extend([p.text.strip() for p in paragraphs if p.text.strip()])
                
                for img in content_div.find_all('img', limit=5):
                    img_url = img.get('src')
                    if img_url:
                        if img_url.startswith('//'):
                            img_url = 'https:' + img_url
                        images.append(img_url)
                
                if content_list:
                    related_keywords = extract_keywords(soup.title.string)
                    break
    
    # If no content found in predefined sources, search using Google
    if not content_list:
        try:
            google_results = google_search(query, num_results=5)  # Fetch top 5 search results
            urls = list(google_results)
            summarized_content = summarize_urls(urls)
            save_to_db(query, summarized_content, [], related_keywords)
            
            return {"content": summarized_content, "images": [], "related_keywords": related_keywords[:3]}
        except Exception as e:
            print(f"Error occurred during Google search: {e}")
            return {"content": "Error occurred during search.", "images": [], "related_keywords": []}
    
    summarized_content = process_content(content_list, query_type, original_query)
    save_to_db(query, summarized_content, images, related_keywords)
    
    return {"content": summarized_content, "images": images, "related_keywords": related_keywords[:3]}
    # Process content based on query type
    summarized_content = process_content(content_list, query_type, original_query)
    save_to_db(query, summarized_content, images, related_keywords)
    
    return {"content": summarized_content, "images": images, "related_keywords": related_keywords[:3]}

# Function to clean text from unnecessary symbols and content
def clean_text(text):
    cleaned_text = re.sub(r'\[.*?\]', '', text)  # Remove text within square brackets
    cleaned_text = re.sub(r'\(.*?\)', '', cleaned_text)  # Remove text within parentheses
    cleaned_text = re.sub(r'\s+', ' ', cleaned_text)  # Remove extra whitespaces
    return cleaned_text.strip()

# Function to summarize content from URLs using LSA
def summarize_urls(urls):
    content_list = []
    
    for url in urls:
        response = requests.get(url)
        if response.status_code == 200:
            soup = BeautifulSoup(response.content, 'html.parser')
            content_div = soup.find('div', class_='content') or soup.find('div', {'id': 'mw-content-text'})
            
            if content_div:
                paragraphs = content_div.find_all('p', limit=25)
                content_list.extend([clean_text(p.text.strip()) for p in paragraphs if p.text.strip()])
    
    full_content = ' '.join(content_list)
    summarized_content = summarize_text(full_content, 10)  # Summarize to 10 sentences
    return summarized_content

# Function to summarize text using LSA
def summarize_text(text, sentence_count):
    parser = PlaintextParser.from_string(text, Tokenizer("english"))
    summarizer = LsaSummarizer()
    summary = summarizer(parser.document, sentence_count)
    return ' '.join([str(sentence) for sentence in summary])

# Function to process content based on query type
def process_content(content_list, query_type, original_query):
    if query_type == "definition":
        return summarize_content(content_list, 1)  # Summarize the first paragraph for definition
    elif query_type == "essay":
        return summarize_content(content_list, 300)  # Summarize up to 300 words for essay
    elif query_type == "analysis":
        return summarize_content(content_list, 25)  # Summarize up to 25 paragraphs for analysis
    elif query_type == "description":
        return summarize_content(content_list, 10)  # Summarize up to 10 paragraphs for description
    else:
        return "Invalid query type."

# Function to join sentences up to a specified number of paragraphs
def summarize_content(content_list, limit):
    joined_content = ' '.join(content_list[:limit])
    return reprocess_text(joined_content, content_list[0] if content_list else "")

# Function to reprocess text to ensure uniqueness
def reprocess_text(text, query):
    language = detect(query)
    if language == 'sw':
        return swahili_transformation(text)
    else:
        return english_transformation(text)

def english_transformation(text):
    sentences = text.split('.')
    random.shuffle(sentences)
    unique_text = '. '.join(sentences)
    unique_text = unique_text.replace(' is ', ' is known to be ').replace(' are ', ' are often ').replace(' was ', ' was found to be ')
    return unique_text

def swahili_transformation(text):
    # Example Swahili transformation rules
    sentences = text.split('.')
    random.shuffle(sentences)
    unique_text = '. '.join(sentences)
    unique_text = unique_text.replace(' ni ', ' inajulikana kuwa ').replace(' wao ni ', ' mara nyingi ni ')
    return unique_text

# Function to save data to MongoDB
def save_to_db(query, content, images, related_keywords):
    db.data.insert_one({
        "query": query,
        "content": content,
        "images": images,
        "related_keywords": related_keywords,
        "timestamp": datetime.now()
    })

# Asynchronous task to update data in the background
@celery.task
def update_data_in_background(query, query_type, original_query):
    result = scrape_information(query, query_type, original_query)
    save_to_db(query, result['content'], result['images'], result['related_keywords'])

# Route for downloading files
@app.route('/download/<file_type>', methods=['POST'])
def download(file_type):
    data = request.form
    query = data.get('query')
    query_type = data.get('query_type')
    content = data.get('content')
    
    if not query or not query_type or not content:
        return render_template('scrpe.html', error="Missing data to generate file.")
    
    result = {"content": content, "images": []}
    
    if file_type == 'pdf':
        create_pdf(query, query_type, result)
        return send_file("output.pdf", as_attachment=True)
    elif file_type == 'ppt':
        create_ppt(query, query_type, result)
        return send_file("Nobestudy.pptx", as_attachment=True)
    elif file_type == 'docx':
        create_docx(query, query_type, result)
        return send_file("Nobestudy.docx", as_attachment=True)
    else:
        return render_template('scrpe.html', error="Invalid file type.")

# Function to create PDF
# Function to create PDF
def create_pdf(query, query_type, result):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    title_text = f"{query_type.capitalize()} of {query}"
    content_text = result['content']
    
    try:
        pdf.cell(200, 10, txt=title_text, ln=True, align='C')
        pdf.multi_cell(0, 10, txt=content_text)
        
        if result['images']:
            for img in result['images']:
                pdf.image(img, w=100)
        
        pdf.output("output.pdf")
    except Exception as e:
        print(f"An error occurred while creating the PDF: {e}")

# Function to create PowerPoint presentation
def create_ppt(query, query_type, result):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"{query_type.capitalize()} of {query}"
    
    left_inch = Inches(1)
    top_inch = Inches(1.2)
    width_inch = Inches(8)
    height_inch = Inches(1.5)
    
    textbox = slide.shapes.add_textbox(left_inch, top_inch, width_inch, height_inch)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = result['content']
    
    if result['images']:
        for img in result['images']:
            slide.shapes.add_picture(img, Inches(1), Inches(2), height=Inches(1.5))
    
    prs.save('Nobestudy.pptx')

# Function to create DOCX document
def create_docx(query, query_type, result):
    doc = Document()
    doc.add_heading(f"{query_type.capitalize()} of {query}", level=1)
    doc.add_paragraph(result['content'])
    
    if result['images']:
        for img in result['images']:
            doc.add_picture(img, width=Inches(2))
    
    doc.save('Nobestudy.docx')

# Run Flask app
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 4000))
    app.run(host="0.0.0.0", port=port)

