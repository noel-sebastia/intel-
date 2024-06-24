import os
import re
import random
from datetime import datetime
from flask import Flask, request, send_file, jsonify, render_template
from bs4 import BeautifulSoup
import requests
from fpdf import FPDF
from pptx import Presentation
from docx import Document
from pptx.util import Inches
from pymongo import MongoClient
from langdetect import detect
from transformers import AutoTokenizer, AutoModelForCausalLM, pipeline
import speech_recognition as sr
import pyttsx3

# Initialize Flask app
app = Flask(__name__)

# Initialize MongoDB
client = MongoClient("mongodb://localhost:27017/")
db = client["scraped_data"]

# Initialize speech recognition and text-to-speech engines
recognizer = sr.Recognizer()
engine = pyttsx3.init()

# Initialize GPT-3 or GPT-Neo model and tokenizer
model_name = "EleutherAI/gpt-neo-125M"  # Adjust model size as needed (e.g., 1.3B for GPT-3, or larger variants)
tokenizer = AutoTokenizer.from_pretrained(model_name)
model = AutoModelForCausalLM.from_pretrained(model_name)
gpt3_pipeline = pipeline("text-generation", model=model, tokenizer=tokenizer)

# List of domains to avoid in content extraction
avoid_domains = ["wikipedia.org", "byjus.com", "wikihow.com", "britannica.com"]

# Route for home page
@app.route('/')
def home():
    return render_template('scrpe.html')

# Route for handling search request
@app.route('/search', methods=['POST'])
def search():
    data = request.get_json()
    query = data.get('query')
    query_type = data.get('query_type')
    
    if not query or not query_type:
        return jsonify({"error": "Please provide both query and query type."}), 400
    
    keywords = extract_keywords(query)
    result = scrape_information(' '.join(keywords), query_type, query)
    
    return jsonify(result)

# Function to extract keywords from a query
def extract_keywords(query):
    return re.findall(r'\b\w+\b', query)

# Function to scrape information from different sources or Google if needed
def scrape_information(query, query_type, original_query):
    sources = [
        f"https://simple.wikipedia.org/wiki/{query}",
        f"https://scholar.google.com/scholar?q={query}",
        f"https://en.wikipedia.org/wiki/{query}",
        f"https://sw.wikipedia.org/wiki/{query}"
    ]
    content_list, images, related_keywords = [], [], []
    
    # Scraping from predefined sources
    for url in sources:
        if any(domain in url for domain in avoid_domains):
            continue
        
        response = requests.get(url)
        if response.status_code == 200:
            soup = BeautifulSoup(response.content, 'html.parser')
            content_div = soup.find('div', class_='content') or soup.find('div', {'id': 'mw-content-text'})
            
            if content_div:
                paragraphs = content_div.find_all('p', limit=25)
                content_list.extend([clean_text(p.text.strip()) for p in paragraphs if p.text.strip()])
                
                for img in content_div.find_all('img', limit=5):
                    img_url = img.get('src')
                    if img_url:
                        if img_url.startswith('//'):
                            img_url = 'https:' + img_url
                        images.append(img_url)
                
                if content_list:
                    related_keywords = extract_keywords(soup.title.string)
                    break
    
    # If no content found in predefined sources, use GPT-3 or GPT-Neo for text generation
    if not content_list:
        try:
            language = detect(original_query)
            gpt3_response = generate_text_with_gpt3(original_query, query_type, language)
            return gpt3_response
        except Exception as e:
            print(f"Error occurred during GPT-3 text generation: {e}")
            return {"content": "Error occurred during search.", "images": [], "related_keywords": []}
    
    summarized_content = process_content(content_list, query_type, original_query)
    save_to_db(query, summarized_content, images, related_keywords)
    
    return {"content": summarized_content, "images": images, "related_keywords": related_keywords[:3]}

# Function to clean text from unnecessary symbols and content
def clean_text(text):
    cleaned_text = re.sub(r'\[.*?\]', '', text)  # Remove text within square brackets
    cleaned_text = re.sub(r'\(.*?\)', '', cleaned_text)  # Remove text within parentheses
    cleaned_text = re.sub(r'\s+', ' ', cleaned_text)  # Remove extra whitespaces
    return cleaned_text.strip()

# Function to process content based on query type
def process_content(content_list, query_type, original_query):
    if query_type == "definition":
        return summarize_content(content_list, 1)  # Summarize the first paragraph for definition
    elif query_type == "essay":
        return summarize_content(content_list, 300)  # Summarize up to 300 words for essay
    elif query_type == "analysis":
        return summarize_content(content_list, 25)  # Summarize up to 25 paragraphs for analysis
    elif query_type == "description":
        return summarize_content(content_list, 10)  # Summarize up to 10 paragraphs for description
    else:
        return "Invalid query type."

# Function to join sentences up to a specified number of paragraphs
def summarize_content(content_list, limit):
    joined_content = ' '.join(content_list[:limit])
    return reprocess_text(joined_content, content_list[0] if content_list else "")

# Function to reprocess text to ensure uniqueness
def reprocess_text(text, query):
    language = detect(query)
    if language == 'sw':
        return swahili_transformation(text)
    else:
        return english_transformation(text)

def english_transformation(text):
    sentences = text.split('.')
    random.shuffle(sentences)
    unique_text = '. '.join(sentences)
    unique_text = unique_text.replace(' is ', ' is known to be ').replace(' are ', ' are often ').replace(' was ', ' was found to be ')
    return unique_text

def swahili_transformation(text):
    # Example Swahili transformation rules
    sentences = text.split('.')
    random.shuffle(sentences)
    unique_text = '. '.join(sentences)
    unique_text = unique_text.replace(' ni ', ' inajulikana kuwa ').replace(' wao ni ', ' mara nyingi ni ')
    return unique_text

# Function to save data to MongoDB
def save_to_db(query, content, images, related_keywords):
    db.data.insert_one({
        "query": query,
        "content": content,
        "images": images,
        "related_keywords": related_keywords,
        "timestamp": datetime.now()
    })

# Function to generate text using GPT-3 or GPT-Neo
def generate_text_with_gpt3(query, query_type, language):
    input_text = f"Generate {query_type} about {query}"
    max_length = get_max_length(query_type)
    
    # Generate text using GPT-3 or GPT-Neo
    generated_content = gpt3_pipeline(input_text, max_length=max_length, num_return_sequences=1)
    generated_text = generated_content[0]['generated_text']
    
    # Extract images from the web related to the generated content
    images = extract_images_from_web(generated_text)
    
    return {"content": generated_text, "images": images, "related_keywords": extract_keywords(generated_text)}

# Function to extract images from web based on generated text
def extract_images_from_web(generated_text):
    images = []
    priority_domains = [
        "https://en.wikipedia.org/",
        "https://google.com/"
    ]
    
    for domain in priority_domains:
        url = domain + generated_text
        response = requests.get(url)
        soup = BeautifulSoup(response.content, 'html.parser')
        images.extend([img.get('src') for img in soup.find_all('img')])
    
    return images

# Function to handle audio conversation
@app.route('/audio-conversation', methods=['POST'])
def audio_conversation():
    # Assuming audio file is sent in request (you can adapt this for streaming audio or live mic input)
    audio_file = request.files['audio']
    
    if not audio_file:
        return jsonify({"error": "No audio file found."}), 400
    
    # Convert audio to text
    with sr.AudioFile(audio_file) as source:
        audio_data = recognizer.record(source)
        user_input = recognizer.recognize_google(audio_data)
    
    # Determine language of user input
    language = detect(user_input)
    
    # Generate response based on user input using GPT-3 or GPT-Neo
    response = generate_text_with_gpt3(user_input, "dialogue", language)
    generated_text = response['content']
    
    # Convert generated text to speech
    engine.say(generated_text)
    engine.runAndWait()
    
    return jsonify({"response": generated_text})

# Function to create PDF
def create_pdf(query, query_type, result):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    title_text = f"{query_type.capitalize()} of {query}"
    content_text = result['content']
    
    try:
        pdf.cell(200, 10, txt=title_text    # in PDF
        , ln=True, align='C')
        pdf.multi_cell(0, 10, txt=content_text)
        
        if result['images']:
            for img in result['images']:
                pdf.image(img, w=100)
        
        pdf.output("output.pdf")
    except Exception as e:
        print(f"An error occurred while creating the PDF: {e}")

# Function to create PowerPoint presentation
def create_ppt(query, query_type, result):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"{query_type.capitalize()} of {query}"
    
    left_inch = Inches(1)
    top_inch = Inches(1.2)
    width_inch = Inches(8)
    height_inch = Inches(1.5)
    
    textbox = slide.shapes.add_textbox(left_inch, top_inch, width_inch, height_inch)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = result['content']
    
    if result['images']:
        for img in result['images']:
            slide.shapes.add_picture(img, Inches(1), Inches(2), height=Inches(1.5))
    
    prs.save('Nobestudy.pptx')

# Function to create DOCX document
def create_docx(query, query_type, result):
    doc = Document()
    doc.add_heading(f"{query_type.capitalize()} of {query}", level=1)
    doc.add_paragraph(result['content'])
    
    if result['images']:
        for img in result['images']:
            doc.add_picture(img, width=Inches(2))
    
    doc.save('Nobestudy.docx')

# Route for downloading files
@app.route('/download/<file_type>', methods=['POST'])
def download(file_type):
    data = request.form
    query = data.get('query')
    query_type = data.get('query_type')
    content = data.get('content')
    
    if not query or not query_type or not content:
        return render_template('scrpe.html', error="Missing data to generate file.")
    
    result = {"content": content, "images": []}
    
    if file_type == 'pdf':
        create_pdf(query, query_type, result)
        return send_file("output.pdf", as_attachment=True)
    elif file_type == 'ppt':
        create_ppt(query, query_type, result)
        return send_file("Nobestudy.pptx", as_attachment=True)
    elif file_type == 'docx':
        create_docx(query, query_type, result)
        return send_file("Nobestudy.docx", as_attachment=True)
    else:
        return render_template('scrpe.html', error="Invalid file type.")

# Run Flask app
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 4000))
    app.run(host="0.0.0.0", port=port)

